# routers/ai_assistant.py
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.orm import Session
import os
from openai import OpenAI
from pydantic import BaseModel
from typing import Optional
from pptx import Presentation
import io
from .. import models
from ..auth import get_db

router = APIRouter(prefix="/ai", tags=["AI Assistant"])


# 初始化AI客户端
def get_ai_client():
    deepseek_api_key = os.getenv("DEEPSEEK_API_KEY", "sk-c641fc0377ea4b51a7c94cd5f6803b0c")
    if not deepseek_api_key:
        raise HTTPException(status_code=500, detail="服务器未配置 AI 模型 API 密钥")

    return OpenAI(
        api_key=deepseek_api_key,
        base_url="https://api.deepseek.com"
    )


# 提取PPTX文字的函数
def extract_pptx_text(file: UploadFile):
    try:
        ppt = Presentation(io.BytesIO(file.file.read()))
        full_text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX文件解析失败: {str(e)}")



class TextQuestionRequest(BaseModel):
    question: str



class PPTXQuestionRequest(BaseModel):
    question: str
    pptx_file: UploadFile


@router.post("/ask")
async def ask_text_question(
        request: TextQuestionRequest,
        db: Session = Depends(get_db)
):
    """处理纯文本提问"""
    question_text = request.question.strip()
    if not question_text:
        raise HTTPException(status_code=400, detail="问题内容不能为空")

    try:
        client = get_ai_client()

        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个乐于助人的 AI 助手。"},
                {"role": "user", "content": question_text}
            ],
            stream=False
        )

        answer_text = response.choices[0].message.content


        record = models.AIRequestRecord(question=question_text, answer=answer_text)
        db.add(record)
        db.commit()

        return {"answer": answer_text}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"服务处理失败: {str(e)}")


@router.post("/ask_with_pptx")
async def ask_with_pptx(
        question: str = Form(...),
        pptx_file: UploadFile = File(...),
        db: Session = Depends(get_db)
):
    """处理带PPTX文件上传的提问"""
    # 验证文件类型
    if not pptx_file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="仅支持PPTX格式文件")

    # 提取文字内容
    try:
        extracted_text = extract_pptx_text(pptx_file)
    except HTTPException as he:
        raise he

    # 合并问题和PPT内容
    final_prompt = f"{question.strip()}\n\n相关PPT内容：\n{extracted_text[:3000]}"

    try:
        client = get_ai_client()

        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个擅长分析PPT内容的AI助手。"},
                {"role": "user", "content": final_prompt}
            ],
            stream=False
        )

        answer_text = response.choices[0].message.content

        # 保存记录（仅保存用户原始问题）
        record = models.AIRequestRecord(question=question, answer=answer_text)
        db.add(record)
        db.commit()

        return {
            "answer": answer_text,
            "ppt_summary": f"已分析{len(extracted_text)}字符的PPT内容"
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"服务处理失败: {str(e)}")